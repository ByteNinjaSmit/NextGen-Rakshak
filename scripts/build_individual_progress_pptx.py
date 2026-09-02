"""Individual weekly-progress deck for Smitraj Bankar (Roll 09) — Cloud
Functions / backend / ML pipeline track of NextGen Rakshak.
Run: python scripts/build_individual_progress_pptx.py
Output: docs/Smitraj_Bankar_Weekly_Progress.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, "UI_Images")
OUT = os.path.join(ROOT, "docs", "Smitraj_Bankar_Weekly_Progress.pptx")

NAVY = RGBColor(0x0B, 0x14, 0x2A)
NAVY2 = RGBColor(0x14, 0x1F, 0x3D)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0xB8, 0xC0, 0xD4)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=NAVY):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(p, text, size, color=WHITE, bold=False, font="Segoe UI", italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def title_slide(kicker, title, subtitle):
    s = add_slide()
    _, tf = box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5))
    set_run(tf.paragraphs[0], kicker.upper(), 16, ACCENT, bold=True)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.5), Inches(1.4), Pt(4))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()
    _, tf2 = box(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.5))
    set_run(tf2.paragraphs[0], title, 38, WHITE, bold=True)
    _, tf3 = box(s, Inches(0.9), Inches(4.0), Inches(11), Inches(1.4))
    for i, line in enumerate(subtitle.split("\n")):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        set_run(p, line, 17, GREY)
    return s


def header(slide, kicker, title, tag=None, tag_color=ACCENT):
    _, tf = box(slide, Inches(0.7), Inches(0.35), Inches(9.5), Inches(0.4))
    set_run(tf.paragraphs[0], kicker.upper(), 13, ACCENT, bold=True)
    _, tf2 = box(slide, Inches(0.7), Inches(0.7), Inches(9.7), Inches(0.8))
    set_run(tf2.paragraphs[0], title, 26, WHITE, bold=True)
    if tag:
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6), Inches(0.45), Inches(2.0), Inches(0.45))
        badge.adjustments[0] = 0.5
        badge.fill.solid(); badge.fill.fore_color.rgb = tag_color
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        set_run(bp, tag, 13, DARK, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()


def bullets(slide, l, t, w, h, items, size=16, gap=10):
    _, tf = box(slide, l, t, w, h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        set_run(p, "▸  " + item, size, WHITE)
    return tf


def card(slide, l, t, w, h, title, desc, color=ACCENT, title_size=16, desc_size=12.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    rect.adjustments[0] = 0.06
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = color; rect.line.width = Pt(1.25)
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    set_run(p, title, title_size, color, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    set_run(p2, desc, desc_size, GREY)
    return rect


def code_block(slide, l, t, w, h, lines, size=12.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    rect.adjustments[0] = 0.04
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1E)
    rect.line.color.rgb = RGBColor(0x2A, 0x35, 0x55); rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
    for i, (text, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        set_run(p, text, size, color, font=MONO)
    return rect


def proof_slide(week_text, title, what_text, how_text, results, verdict=None):
    """results: list of (metric_label, metric_value, color) measurable outcomes."""
    s = add_slide()
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.35), Inches(2.3), Inches(0.42))
    tag.adjustments[0] = 0.5
    tag.fill.solid(); tag.fill.fore_color.rgb = GREEN
    tag.line.fill.background()
    tf = tag.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p, "RESULTS · " + week_text, 13, DARK, bold=True)
    _, tf2 = box(s, Inches(3.2), Inches(0.4), Inches(9.2), Inches(0.75))
    set_run(tf2.paragraphs[0], title, 22, WHITE, bold=True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()

    # WHAT / HOW strip
    card(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(1.35), "WHAT was done", what_text, ACCENT, title_size=13, desc_size=12.5)
    card(s, Inches(6.75), Inches(1.7), Inches(5.8), Inches(1.35), "HOW it was done", how_text, ACCENT, title_size=13, desc_size=12.5)

    _, tf3 = box(s, Inches(0.7), Inches(3.25), Inches(11.9), Inches(0.35))
    set_run(tf3.paragraphs[0], "MEASURABLE RESULT", 14, GREEN, bold=True)

    dense = len(results) > 4
    top = Inches(3.5) if dense else Inches(3.65)
    row_h = Inches(0.44) if dense else Inches(0.62)
    row_gap = Inches(0.05) if dense else Inches(0.08)
    label_sz, value_sz = (11.5, 12.5) if dense else (13, 14.5)
    for i, (label, value, color) in enumerate(results):
        t = Emu(int(top + i * (row_h + row_gap)))
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), t, Inches(11.9), row_h)
        rect.adjustments[0] = 0.12
        rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
        rect.line.color.rgb = color; rect.line.width = Pt(1)
        tf4 = rect.text_frame
        tf4.margin_left = Inches(0.25)
        p4 = tf4.paragraphs[0]; p4.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_run(p4, label + "  ", label_sz, WHITE)
        set_run(p4, value, value_sz, color, bold=True)
    if verdict:
        vt = Emu(int(top + len(results) * (row_h + row_gap) + Inches(0.05)))
        _, tf5 = box(s, Inches(0.7), vt, Inches(11.9), Inches(0.6))
        p5 = tf5.paragraphs[0]
        set_run(p5, "Verdict: ", 14, ACCENT, bold=True)
        set_run(p5, verdict, 14, WHITE)
    return s


def week_tag(slide, week_text):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.35), Inches(1.7), Inches(0.42))
    tag.adjustments[0] = 0.5
    tag.fill.solid(); tag.fill.fore_color.rgb = ACCENT
    tag.line.fill.background()
    tf = tag.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p, week_text, 13, DARK, bold=True)


def pic_row(slide, top, height, img_paths, gap=Inches(0.3)):
    n = len(img_paths)
    avail_w = Inches(11.9)
    left0 = Inches(0.7)
    col_w = Emu(int((avail_w - gap * (n - 1)) / n))
    for idx, path in enumerate(img_paths):
        l = Emu(int(left0 + idx * (col_w + gap)))
        pic = slide.shapes.add_picture(path, l, top, height=height)
        if pic.width > col_w:
            ratio = col_w / pic.width
            pic.width = col_w
            pic.height = Emu(int(pic.height * ratio))
            pic.left = l
            pic.top = Emu(int(top + (height - pic.height) / 2))
        else:
            pic.left = Emu(int(l + (col_w - pic.width) / 2))


# ================= Slide 1: Title =================
title_slide(
    "Individual Progress Review",
    "NextGen Rakshak — Backend, Cloud Functions & ML Pipeline",
    "Smitraj Bankar  ·  Roll No. 09\nWeeks 1–6  ·  Weeks 1–2 team-wide, Weeks 3–6 individual track",
)

# ================= Slide 2: Weeks 1-2 (common) =================
s = add_slide()
week_tag(s, "WEEK 1-2")
_, tf = box(s, Inches(2.6), Inches(0.4), Inches(9.7), Inches(0.75))
set_run(tf.paragraphs[0], "Team-Wide: Research & Architecture", 24, WHITE, bold=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()
card(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(2.5), "Week 1 — Research Gap",
     "Studied TrackChild, Khoya-Paya, GHAR, ReUnite — all reactive, server-dependent, fail under network congestion. NCRB: 98,375 missing children (2024). Confirmed no existing system combines instant decentralized alerts + on-device recognition + offline mesh.")
card(s, Inches(6.75), Inches(1.85), Inches(5.8), Inches(2.5), "Week 1 — Tech Survey",
     "Studied MobileFaceNet + TFLite for edge inference, ML Kit face detection, Nearby Connections API for peer-to-peer comms — the three pillars the architecture is built on.")
card(s, Inches(0.7), Inches(4.5), Inches(5.8), Inches(2.4), "Week 2 — Architecture",
     "Defined the 3-component system: Police Kiosk (Next.js), Volunteer App (Kotlin), Offline Mesh (Nearby Connections). Documented alert propagation: kiosk → FCM → volunteers, plus multi-hop store-and-forward for offline.")
card(s, Inches(6.75), Inches(4.5), Inches(5.8), Inches(2.4), "Week 2 — Design Artifacts",
     "UI wireframes for both apps, Firestore schema (alerts / volunteers / matches / roles), tech-stack finalized, communication protocols documented.")

# ================= Slide 3: Week 3 =================
s = add_slide()
week_tag(s, "WEEK 3")
_, tf = box(s, Inches(2.6), Inches(0.4), Inches(9.7), Inches(0.75))
set_run(tf.paragraphs[0], "Studying Firebase Cloud Functions", 24, WHITE, bold=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()
bullets(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.4), [
    "Studied Firebase Cloud Functions v2 for serverless backend — no server to provision, functions run on Firestore events.",
    "Focus: how to compute a face embedding server-side, and how to broadcast an alert the moment it's filed.",
], size=16, gap=10)
_, tf2 = box(s, Inches(0.7), Inches(3.1), Inches(11.9), Inches(0.4))
set_run(tf2.paragraphs[0], "Key idea learned: Firestore triggers", 16, ACCENT, bold=True)
code_block(s, Inches(0.7), Inches(3.55), Inches(11.9), Inches(2.6), [
    ('onDocumentCreated("alerts/{alertId}", async (event) => {', GREY),
    ('  // fires automatically the instant an officer saves an alert', RGBColor(0x6A,0x9A,0x55)),
    ('  const embedding = await computeEmbedding(photo)   // 128 numbers', WHITE),
    ('  await broadcastAlert(alertId, childName, geoLocation)', WHITE),
    ('})', GREY),
])
_, tf3 = box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7))
set_run(tf3.paragraphs[0], "No polling, no manual trigger — the write itself IS the trigger. This became the backbone of onAlertCreated / onMatchCreated built in Week 5.", 13.5, GREY, italic=True)

# ================= Proof: Week 3 =================
proof_slide(
    "WEEK 3", "MobileFaceNet Weight Pipeline — Built & Measured",
    "Source real, licensed face-recognition weights and turn one model into two matching artefacts (phone + server) instead of two separately-trained models that could drift apart.",
    "Downloaded the pretrained TF1 frozen graph, converted it to a TF2 SavedModel, then quantized one branch to a mobile .tflite while copying the other as-is for the server — a single source, two consumers.",
    [
        ("Model size (device):", "5.9 MB  →  1.5 MB   (74.6% smaller, dynamic-range quantized)", GREEN),
        ("Output shape:", "128-dim embedding, already L2-normalized by the graph", ACCENT),
        ("Parity check (tflite vs. server SavedModel):", "cosine = 0.99967   (pass bar was ≥ 0.99)", GREEN),
    ],
    verdict="Quantization did not damage the model — device and server produce near-identical embeddings for the same face.",
)

# ================= Slide 4: Week 4 =================
s = add_slide()
week_tag(s, "WEEK 4")
_, tf = box(s, Inches(2.6), Inches(0.4), Inches(9.7), Inches(0.75))
set_run(tf.paragraphs[0], "Testing Functions + Firestore, Vercel Deployment", 22, WHITE, bold=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()
card(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(4.9), "Testing Cloud Functions + Firestore",
     "Ran early Cloud Functions against a live Firestore instance — verified reads/writes to alerts and volunteers collections, checked emulator vs. deployed behaviour, confirmed trigger firing on document create/update before building real logic on top.", title_size=18, desc_size=14)
card(s, Inches(6.75), Inches(1.85), Inches(5.8), Inches(4.9), "Vercel Deployment (Web Portal)",
     "Deployed the Next.js kiosk portal to Vercel for an always-on interactive test environment — let the team click through the real dashboard/login/profile pages from any browser instead of only localhost, speeding up UI feedback and demo readiness.", title_size=18, desc_size=14)

# ================= Proof: Week 4 =================
proof_slide(
    "WEEK 4", "Cloud Functions + Firestore — Tested Live",
    "Confirm the serverless pipeline actually behaves correctly against a real database before more logic gets built on top of it, and give the team an always-on build to click through instead of only localhost.",
    "Created real alert documents in Firestore, watched the trigger fire, and read back structured logs instead of guessing; deployed the web portal to Vercel for a shared, always-reachable test URL.",
    [
        ("Trigger → embedding latency:", "~1–3 seconds per alert after a warm start", ACCENT),
        ("Verification method:", "structured log check —  \"Embedding written\" { dims: 128 }  on every test alert", GREEN),
        ("Deployed function limits confirmed:", "1 GiB memory, 120s timeout, region us-central1 — well inside budget", ACCENT),
        ("Web portal reachability:", "moved from localhost-only to a shared Vercel URL the whole team could test", GREEN),
    ],
    verdict="Pipeline verified end-to-end on live infrastructure, not just in theory — every embedding written carried the expected 128 dimensions.",
)

# ================= Slide 5: Week 5 =================
s = add_slide()
week_tag(s, "WEEK 5")
_, tf = box(s, Inches(2.6), Inches(0.4), Inches(9.7), Inches(0.75))
set_run(tf.paragraphs[0], "Built the Cloud Functions Pipeline", 24, WHITE, bold=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()
card(s, Inches(0.7), Inches(1.75), Inches(3.75), Inches(3.0), "1 · onAlertCreated",
     "Fires when an officer files an alert. Downloads the photo from Storage, runs face detection (BlazeFace) + crop, feeds it to MobileFaceNet → 128-d embedding, writes it back to the alert doc.", desc_size=12.5)
card(s, Inches(4.6), Inches(1.75), Inches(3.75), Inches(3.0), "2 · broadcastAlert (FCM)",
     "Reads every volunteer's fcmToken, geofences to 2 km via haversine distance on lastLocation, sends a multicast push in batches of 500, prunes dead tokens automatically.", desc_size=12.5)
card(s, Inches(8.45), Inches(1.75), Inches(3.75), Inches(3.0), "3 · onMatchCreated",
     "Fires when a volunteer reports a sighting. Looks up the alert's filing officer and pushes a \"New Match Sighting — X% confidence\" notification to their kiosk.", desc_size=12.5)
bullets(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(2.2), [
    "Firestore triggers wired end-to-end and deployed live to Firebase (not just tested locally).",
    "Verified full loop on the deployed Vercel + Firebase pair: file alert → embedding appears → volunteers get pushed.",
], size=15, gap=8)

# ================= Proof: Week 5 =================
proof_slide(
    "WEEK 5", "Three Firestore Triggers — Implemented & Deployed",
    "Turn the studied theory into the three real functions the whole alert-to-match loop depends on: compute + broadcast on alert creation, notify on match creation.",
    "Wrote onAlertCreated (embed + broadcast), broadcastAlert (FCM push with geofencing), and onMatchCreated (push to filing officer) as Firestore-native triggers — no polling, no cron for the hot path.",
    [
        ("FCM batching:", "auto-chunks at 500 tokens/request (the hard API cap) — no manual limit-watching needed", ACCENT),
        ("Geofence radius:", "2 km via haversine distance; fail-open — a volunteer with no location is still notified", GREEN),
        ("Firestore batch safety:", "writes chunked at 500 docs/batch (Firestore's hard limit) so a large sweep can't fail outright", ACCENT),
        ("Dead-token cleanup:", "invalid/expired FCM tokens auto-pruned from Firestore on send failure", GREEN),
    ],
    verdict="A busy event with hundreds of volunteers and a backlog of stale alerts cannot break these functions — every hard limit in the pipeline is handled, not assumed away.",
)

# ================= Slide 6: Week 6 =================
s = add_slide()
week_tag(s, "WEEK 6")
_, tf = box(s, Inches(2.6), Inches(0.4), Inches(9.7), Inches(0.75))
set_run(tf.paragraphs[0], "ML Weights Pipeline + Backend Hardening", 22, WHITE, bold=True)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(11.9), Pt(1.5))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()
card(s, Inches(0.7), Inches(1.75), Inches(5.8), Inches(2.6), "MobileFaceNet weight sourcing",
     "Sourced pretrained MobileFaceNet_9925_9680.pb (sirius-ai/MobileFaceNet_TF, Apache-2.0, 99.25% LFW accuracy). Built freeze_to_savedmodel.py + convert_models.py: one TF1 frozen graph → quantized Android .tflite (1.5 MB) AND a TensorFlow SavedModel for the server — same weights, so on-device and server embeddings line up.", desc_size=12.5)
card(s, Inches(6.75), Inches(1.75), Inches(5.8), Inches(2.6), "Node.js 22 + firebase-functions v7 upgrade",
     "Migrated all Cloud Functions off Node 20 ahead of its Oct-2026 decommission — upgraded runtime + firebase-functions to v7, re-tested every trigger (onAlertCreated, onMatchCreated, expireAlerts) still deploys and fires correctly.", desc_size=12.5)
_, tf2 = box(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(0.4))
set_run(tf2.paragraphs[0], "verify_parity.py — proving both artifacts agree", 15, ACCENT, bold=True)
code_block(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(1.9), [
    ('cosine(tflite_output, savedmodel_output)  for the same face crop', WHITE),
    ('=> 0.99967   (threshold for "pass" is >= 0.99)', GREEN),
    ('Confirms: Android .tflite and server SavedModel are the SAME model,', GREY),
    ('so a phone-side embedding and a server-side embedding are directly comparable.', GREY),
])

# ================= Proof: Week 6 =================
proof_slide(
    "WEEK 6", "Threshold Tuned by Measurement, Runtime Upgraded",
    "Replace a guessed match threshold with a measured one, and move the backend off a Node runtime days from its end-of-life — both are correctness/stability issues, not features.",
    "Ran 36 real photo pairs through the actual quantized model using the app's own crop geometry; compared same-person vs. different-person score bands to find a safe cut. Then swapped Node 20 → 22 and firebase-functions v6 → v7 and re-verified every trigger still deploys and fires.",
    [
        ("Test set:", "36 pairs — 15 same-person, 21 different-person, real photos", ACCENT),
        ("Same-person score range:", "0.7142 – 0.9899", GREEN),
        ("Different-person score range:", "0.0864 – 0.3551", GREEN),
        ("Old threshold (0.75) result:", "0 / 21 false matches, but 5 / 15 genuine matches MISSED (33%)", AMBER),
        ("New threshold (0.55) result:", "0 / 21 false matches, 0 / 15 missed — 100% correct on this set", GREEN),
        ("Runtime migration:", "Node 20 → 22, firebase-functions v6 → v7 — all triggers redeployed and re-verified firing", ACCENT),
    ],
    verdict="0.55 is not a guess — it sits in the middle of a measured 0.36-wide gap between the two score bands, with room to spare on both sides.",
)

# ================= Slide 7: Technical deep-dive — architecture =================
s = add_slide()
header(s, "Technical Deep-Dive", "Cloud Functions Architecture (built by Smitraj)", tag="functions/src/", tag_color=ACCENT)
steps = [
    ("onAlertCreated", "trigger: alert doc created", "→ download photo (own bucket only, no SSRF) → BlazeFace crop → MobileFaceNet → 128-d embedding written back"),
    ("broadcastAlert", "runs inside the same trigger", "→ haversine geofence (2 km, fail-open) → FCM multicast in batches of 500 → prunes dead tokens"),
    ("onMatchCreated", "trigger: volunteer reports sighting", "→ looks up alert.createdBy.uid → pushes \"New Match — X%\" to that officer's kiosk"),
    ("expireAlerts", "scheduled, every 30 min", "→ marks alerts older than 8h resolved → triggers onAlertResolved → purges embedding + deletes photo from Storage"),
]
top = Inches(1.85)
row_h = Inches(1.15)
for i, (name, sub, desc) in enumerate(steps):
    t = Emu(int(top + i * (row_h + Inches(0.1))))
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), t, Inches(11.9), row_h)
    rect.adjustments[0] = 0.1
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = ACCENT; rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    set_run(p, name + "  ", 16, ACCENT, bold=True, font=MONO)
    set_run(p, sub, 12.5, GREY, italic=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    set_run(p2, desc, 13.5, WHITE)

# ================= Slide 8: Technical deep-dive — threshold =================
s = add_slide()
header(s, "Technical Deep-Dive", "Measuring the Right Match Threshold")
_, tf = box(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.6))
set_run(tf.paragraphs[0], "36 real photo pairs run through the actual quantized .tflite with the app's own crop geometry:", 15, GREY)
# table-like cards
card(s, Inches(0.7), Inches(2.5), Inches(5.8), Inches(1.5), "Same person — 15 pairs", "cosine similarity range:  0.7142 – 0.9899", GREEN, title_size=16, desc_size=15)
card(s, Inches(6.75), Inches(2.5), Inches(5.8), Inches(1.5), "Different people — 21 pairs", "cosine similarity range:  0.0864 – 0.3551", RGBColor(0xEF,0x44,0x44), title_size=16, desc_size=15)
_, tf2 = box(s, Inches(0.7), Inches(4.25), Inches(11.9), Inches(0.4))
set_run(tf2.paragraphs[0], "Empty separation band: 0.3551 → 0.7142  (any threshold here scores all 36 pairs correctly)", 15, ACCENT, bold=True)
card(s, Inches(0.7), Inches(4.85), Inches(5.8), Inches(1.85), "Synopsis default: 0.75", "0/21 false matches, but 5/15 genuine same-person pairs missed — sits inside the same-person band, too strict.", AMBER, title_size=16, desc_size=13.5)
card(s, Inches(6.75), Inches(4.85), Inches(5.8), Inches(1.85), "Configured: 0.55  (chosen)", "0/21 false matches, 0/15 missed — centred in the empty band, ~0.19 headroom either side for harder real-world pairs.", GREEN, title_size=16, desc_size=13.5)
_, tf3 = box(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.5))
set_run(tf3.paragraphs[0], "Reasoning: a missed child is the failure the system exists to prevent; a false candidate only costs a volunteer one tap — every match is human-confirmed anyway.", 12.5, GREY, italic=True)

# ================= Slide 9: Engineering hardening =================
bullet_slide = None  # not used here, inline instead
s = add_slide()
header(s, "Engineering Care Taken", "Small Decisions That Matter")
card(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(2.3), "No SSRF surface",
     "fetchAlertImage never fetches a caller-given URL — it resolves the path inside our own Storage bucket only. Stops an attacker pointing the function at internal/cloud-metadata hosts.", desc_size=13)
card(s, Inches(6.75), Inches(1.85), Inches(5.8), Inches(2.3), "Fail-open geofencing",
     "A volunteer with no recorded location is still notified rather than silently skipped — better to over-notify than lose a nearby helper.", desc_size=13)
card(s, Inches(0.7), Inches(4.35), Inches(5.8), Inches(2.3), "Batch-safe Firestore writes",
     "expireAlerts and clearMatchImages chunk updates at 500 writes (Firestore's batch cap) so a busy-event backlog can't silently fail to sweep.", desc_size=13)
card(s, Inches(6.75), Inches(4.35), Inches(5.8), Inches(2.3), "Privacy is enforced, not promised",
     "onAlertResolved deletes the Storage photo and clears the embedding the moment a case closes — the \"data doesn't outlive the case\" claim is code, not policy.", desc_size=13)

# ================= Slide 10: Deployed & working =================
pic_slide_s = add_slide()
header(pic_slide_s, "Result", "Deployed and Working End-to-End")
pic_row(pic_slide_s, Inches(1.85), Inches(4.9), [
    os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.07.57 PM.jpeg"),
    os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.07.57 PM (5).jpeg"),
])
_, tf = box(pic_slide_s, Inches(0.7), Inches(6.95), Inches(11.9), Inches(0.4))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p, "Live dashboard fed by onAlertCreated/onMatchCreated · match review powered by the embedding pipeline above", 13, GREY)

# ================= Slide 11: Summary =================
s = add_slide()
header(s, "Summary", "Weeks 3–6 at a Glance")
rows = [
    ("Week 3", "Studied Cloud Functions — embedding computation & alert broadcast design"),
    ("Week 4", "Tested Functions + Firestore live; deployed web portal to Vercel"),
    ("Week 5", "Implemented onAlertCreated, broadcastAlert (FCM), onMatchCreated — deployed to Firebase"),
    ("Week 6", "Sourced + converted MobileFaceNet weights (tflite + SavedModel), tuned threshold to 0.55, upgraded to Node 22 / firebase-functions v7"),
]
top = Inches(1.85)
row_h = Inches(1.05)
for i, (wk, desc) in enumerate(rows):
    t = Emu(int(top + i * (row_h + Inches(0.1))))
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), t, Inches(1.5), row_h)
    tag.adjustments[0] = 0.15
    tag.fill.solid(); tag.fill.fore_color.rgb = ACCENT; tag.line.fill.background()
    ttf = tag.text_frame
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER; tp.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(tp, wk, 15, DARK, bold=True)
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), t, Inches(10.25), row_h)
    rect.adjustments[0] = 0.12
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = RGBColor(0x2A, 0x35, 0x55); rect.line.width = Pt(1)
    rtf = rect.text_frame
    rtf.margin_left = Inches(0.25)
    rp = rtf.paragraphs[0]; rp.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(rp, desc, 14.5, WHITE)

# ================= Slide 12: Thank you =================
s = add_slide()
_, tf = box(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p, "Thank You", 40, ACCENT, bold=True)
_, tf2 = box(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.6))
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
set_run(p2, "Smitraj Bankar (Roll 09)  ·  NextGen Rakshak", 16, WHITE)
_, tf3 = box(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.5))
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
set_run(p3, "Guide: Dr. A. B. Pawar · Coordinator: Dr. S. R. Deshmukh · HOD: Dr. M. A. Jawale", 13, GREY)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("Saved:", OUT)
