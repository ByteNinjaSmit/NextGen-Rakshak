"""Generate the guide-review progress deck for NextGen Rakshak.
Run: python scripts/build_progress_pptx.py
Output: docs/NextGen_Rakshak_Progress.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, "UI_Images")
OUT = os.path.join(ROOT, "docs", "NextGen_Rakshak_Progress.pptx")

NAVY = RGBColor(0x0B, 0x14, 0x2A)
NAVY2 = RGBColor(0x14, 0x1F, 0x3D)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0xB8, 0xC0, 0xD4)
DARK = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=NAVY):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(p, text, size, color=WHITE, bold=False, font="Segoe UI"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return r


def title_slide(title, subtitle, kicker=None):
    s = add_slide()
    if kicker:
        _, tf = box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.6))
        p = tf.paragraphs[0]
        set_run(p, kicker.upper(), 16, ACCENT, bold=True)
    _, tf = box(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.6))
    p = tf.paragraphs[0]
    set_run(p, title, 40, WHITE, bold=True)
    if subtitle:
        _, tf2 = box(s, Inches(0.9), Inches(4.0), Inches(11), Inches(1.2))
        p2 = tf2.paragraphs[0]
        set_run(p2, subtitle, 18, GREY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.55), Inches(1.4), Pt(4))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()
    return s


def header(slide, kicker, title):
    _, tf = box(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.4))
    p = tf.paragraphs[0]
    set_run(p, kicker.upper(), 13, ACCENT, bold=True)
    _, tf2 = box(slide, Inches(0.7), Inches(0.75), Inches(11.9), Inches(0.8))
    p2 = tf2.paragraphs[0]
    set_run(p2, title, 28, WHITE, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.5), Inches(11.9), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x35, 0x55); ln.line.fill.background()


def bullet_slide(kicker, title, bullets, bullet_color=None):
    s = add_slide()
    header(s, kicker, title)
    _, tf = box(s, Inches(0.9), Inches(1.9), Inches(11.3), Inches(5.2))
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        if isinstance(b, tuple):
            head, desc = b
            set_run(p, "▸  " + head, 20, ACCENT if not bullet_color else bullet_color, bold=True)
            if desc:
                p2 = tf.add_paragraph()
                p2.space_after = Pt(6)
                set_run(p2, "     " + desc, 15, GREY)
        else:
            set_run(p, "•  " + b, 18, WHITE)
    return s


def card(slide, l, t, w, h, title, desc, color=ACCENT):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    rect.adjustments[0] = 0.06
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = color; rect.line.width = Pt(1.25)
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    set_run(p, title, 17, color, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    set_run(p2, desc, 13, GREY)
    return rect


def pic_slide(kicker, title, img_paths, caption=None, cols=1):
    s = add_slide()
    header(s, kicker, title)
    n = len(img_paths)
    top = Inches(1.85)
    avail_w = Inches(11.9)
    avail_h = Inches(5.1)
    gap = Inches(0.3)
    left0 = Inches(0.7)
    col_w = Emu(int((avail_w - gap * (cols - 1)) / cols))
    for idx, path in enumerate(img_paths):
        col = idx % cols
        row = idx // cols
        rows = -(-n // cols)
        row_h = Emu(int((avail_h - gap * (rows - 1)) / rows))
        l = Emu(int(left0 + col * (col_w + gap)))
        t = Emu(int(top + row * (row_h + gap)))
        pic = s.shapes.add_picture(path, l, t, height=row_h)
        if pic.width > col_w:
            ratio = col_w / pic.width
            pic.width = col_w
            pic.height = Emu(int(pic.height * ratio))
            pic.left = l
            pic.top = Emu(int(t + (row_h - pic.height) / 2))
        else:
            pic.left = Emu(int(l + (col_w - pic.width) / 2))
    if caption:
        _, tf = box(s, Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.4))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, caption, 13, GREY)
    return s


# ---------- Slide 1: Title ----------
s = title_slide(
    "NextGen Rakshak",
    "Smart Edge-AI System to Reunite Lost Children with Their Families\nat Crowded Public Events — Final Year Capstone Project",
    kicker="Progress Review",
)
_, tf = box(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.6))
p = tf.paragraphs[0]
set_run(p, "Group 6  |  Guide: Dr. A. B. Pawar", 14, GREY)

# ---------- Slide 2: Problem ----------
s = add_slide()
header(s, "The Problem", "Children Get Lost. The First Hour Decides Everything.")
card(s, Inches(0.7), Inches(2.0), Inches(3.75), Inches(4.5),
     "Golden Hour", "The first 60–90 minutes after a child goes missing matter most — after that, the chances of a quick, safe recovery drop fast.")
card(s, Inches(4.75), Inches(2.0), Inches(3.75), Inches(4.5),
     "Networks Jam Up", "Lakhs of people at one venue means mobile networks choke exactly when they're needed most.")
card(s, Inches(8.8), Inches(2.0), Inches(3.75), Inches(4.5),
     "Manual Search Is Slow", "Announcements and printed photos depend on someone noticing and remembering a face — too slow at scale.")

# ---------- Slide 3: Our Idea ----------
s = add_slide()
header(s, "Our Solution", "Put the Search Party in Everyone's Pocket")
_, tf = box(s, Inches(0.9), Inches(2.0), Inches(11.3), Inches(1.3))
p = tf.paragraphs[0]
set_run(p, "Every volunteer's phone quietly checks nearby faces against the missing child's photo\n"
           "using AI that runs on the phone itself.", 22, WHITE)
card(s, Inches(0.7), Inches(3.5), Inches(3.75), Inches(3.2), "No Internet Needed",
     "Phones relay alerts to each other directly (Bluetooth / Wi-Fi), so the search keeps going even offline.", GREEN)
card(s, Inches(4.75), Inches(3.5), Inches(3.75), Inches(3.2), "Privacy by Design",
     "No face photo or video is ever uploaded. Matching happens entirely on the volunteer's device.", GREEN)
card(s, Inches(8.8), Inches(3.5), Inches(3.75), Inches(3.2), "Human Confirms",
     "AI only suggests a possible match — a real person always confirms before anyone is dispatched.", GREEN)

# ---------- Slide 4: How it works flow ----------
s = add_slide()
header(s, "How It Works", "From Alert to Reunion — Step by Step")
steps = [
    "Officer files an alert at the kiosk (photo + child details)",
    "Cloud turns the photo into a small face \"fingerprint\" (just numbers)",
    "Alert reaches volunteers nearby via push notification — and via phone-to-phone relay if offline",
    "Volunteer scans the crowd; matching runs on their phone, live",
    "On a likely match, the app asks the volunteer to visually confirm",
    "Confirmed sighting (location + time) appears instantly on the officer's dashboard for dispatch",
]
top = Inches(1.85)
step_h = Inches(0.82)
for i, step in enumerate(steps):
    t = Emu(int(top + i * (step_h + Inches(0.08))))
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), t, Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
    ctf = circ.text_frame
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    set_run(cp, str(i + 1), 18, WHITE, bold=True)
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.55), t, Inches(10.9), Inches(0.6))
    rect.adjustments[0] = 0.15
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = RGBColor(0x2A, 0x35, 0x55); rect.line.width = Pt(1)
    rtf = rect.text_frame
    rtf.margin_left = Inches(0.25)
    rp = rtf.paragraphs[0]
    rp.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(rp, step, 15, WHITE)

# ---------- Slide 5: Web portal screenshots ----------
pic_slide(
    "Built & Working — Police Web Portal",
    "Officer Dashboard, New Alert, Live Matches",
    [
        os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.07.57 PM.jpeg"),
        os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.07.57 PM (1).jpeg"),
        os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.07.57 PM (5).jpeg"),
    ],
    caption="Left → right: live dashboard, filing a new alert, reviewing a volunteer's sighting",
    cols=3,
)

# ---------- Slide 6: Mobile app screenshots ----------
pic_slide(
    "Built & Working — Volunteer Mobile App",
    "Scan the Crowd, Confirm a Match, Track Status",
    [
        os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.09.18 PM.jpeg"),
        os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.09.17 PM (1).jpeg"),
        os.path.join(IMG, "WhatsApp Image 2026-08-26 at 6.09.19 PM.jpeg"),
    ],
    caption="Left → right: home with active alert, live match found on-device, volunteer's match history",
    cols=3,
)

# ---------- Slide 7: Progress table ----------
s = add_slide()
header(s, "Progress", "What's Done vs. What's Left")
rows = [
    ("Police Web Portal", "Done", GREEN),
    ("Volunteer Mobile App", "Done", GREEN),
    ("On-Device AI Face Matching", "Done", GREEN),
    ("Cloud Backend (alerts, push, expiry)", "Done", GREEN),
    ("Offline Phone-to-Phone Relay", "Done", GREEN),
    ("Raspberry Pi Gate Camera (optional)", "Not started — stretch goal", RGBColor(0xF5, 0x9E, 0x0B)),
]
top = Inches(1.9)
row_h = Inches(0.78)
for i, (name, status, color) in enumerate(rows):
    t = Emu(int(top + i * (row_h + Inches(0.06))))
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), t, Inches(11.9), row_h)
    rect.adjustments[0] = 0.12
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = RGBColor(0x2A, 0x35, 0x55); rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(p, name, 17, WHITE)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Emu(int(t + row_h/2 - Inches(0.19))), Inches(3.1), Inches(0.38))
    badge.adjustments[0] = 0.5
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    btf = badge.text_frame
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    set_run(bp, status, 12, DARK, bold=True)

# ---------- Slide 8: Why different ----------
bullet_slide(
    "Why It Stands Out",
    "Not Just Another Missing-Person App",
    [
        ("Privacy-first", "No bystander's face, photo, or video ever leaves the phone — only numbers (a \"fingerprint\") are used."),
        ("Works without internet", "Most systems die when networks jam; ours keeps relaying alerts phone-to-phone."),
        ("Human always decides", "AI narrows down candidates; a volunteer visually confirms before dispatch — no automated \"found\" decisions."),
    ],
)

# ---------- Slide 9: Next steps ----------
bullet_slide(
    "Next Steps",
    "What's Left Before Final Submission",
    [
        "(Optional) Raspberry Pi fixed-camera node at entry/exit gates",
        "More field-style testing at a mock event",
        "Polish UI details and rehearse the live demo",
    ],
)

# ---------- Slide 10: Team / thank you ----------
s = add_slide()
header(s, "Team", "Group 6")
names = [
    ("09", "Bankar Smitraj Dinkar"),
    ("11", "Bhakare Tanishka Sharad"),
    ("34", "Dhadge Vedant Sanjay"),
    ("94", "Narkhede Atharva Anantkumar"),
]
top = Inches(2.2)
for i, (roll, name) in enumerate(names):
    t = Emu(int(top + i * Inches(0.9)))
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), t, Inches(8.3), Inches(0.72))
    rect.adjustments[0] = 0.2
    rect.fill.solid(); rect.fill.fore_color.rgb = NAVY2
    rect.line.color.rgb = ACCENT; rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    p.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(p, f"{roll}   {name}", 18, WHITE)
_, tf = box(s, Inches(0.9), Inches(6.2), Inches(11.3), Inches(0.9))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Guide: Dr. A. B. Pawar   |   Coordinator: Dr. S. R. Deshmukh   |   HOD: Dr. M. A. Jawale", 14, GREY)
_, tf2 = box(s, Inches(0.9), Inches(6.8), Inches(11.3), Inches(0.6))
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
set_run(p2, "Thank You", 22, ACCENT, bold=True)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("Saved:", OUT)
